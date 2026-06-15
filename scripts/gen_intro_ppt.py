#!/usr/bin/env python3
"""生成 AgenticDataHub 产品介绍 PPT → docs/AgenticDataHub-产品介绍.pptx

用法：.venv/bin/pip install python-pptx && .venv/bin/python scripts/gen_intro_ppt.py
内容是产品介绍的事实摘要，改了文案重跑即可重新生成。
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "AgenticDataHub-产品介绍.pptx"

BRAND = RGBColor(0x2E, 0x7D, 0x5B)       # 品牌深绿
BRAND_LT = RGBColor(0x52, 0xBD, 0x94)    # 品牌绿
INK = RGBColor(0x1F, 0x29, 0x37)         # 深墨
GRAY = RGBColor(0x5B, 0x66, 0x70)
BG = RGBColor(0xF4, 0xF7, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()


def box(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(4)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
        r.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    _fill(s, color); return s


def bg(slide, color=BG):
    rect(slide, 0, 0, SW, SH, color)


def title_slide(title, subtitle, tag):
    s = prs.slides.add_slide(BLANK)
    bg(s, INK)
    rect(s, 0, int(SH * 0.62), SW, Emu(26000), BRAND_LT)
    box(s, Emu(700000), int(SH * 0.30), int(SW * 0.85), Emu(1400000), title, 54, WHITE, bold=True)
    box(s, Emu(720000), int(SH * 0.50), int(SW * 0.85), Emu(700000), subtitle, 24, BRAND_LT)
    box(s, Emu(720000), int(SH * 0.70), int(SW * 0.85), Emu(900000), tag, 16, RGBColor(0xC8, 0xD2, 0xCE))
    return s


def content_slide(no, title, bullets, foot=None):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    rect(s, Emu(560000), Emu(430000), Emu(150000), Emu(560000), BRAND_LT)
    box(s, Emu(800000), Emu(400000), int(SW * 0.82), Emu(650000), title, 32, INK, bold=True)
    box(s, Emu(560000), Emu(380000), Emu(220000), Emu(300000), f"{no:02d}", 16, BRAND, bold=True)
    y = Emu(1300000)
    for head, sub in bullets:
        rect(s, Emu(800000), y + Emu(40000), Emu(150000), Emu(150000), BRAND_LT)
        box(s, Emu(1050000), y - Emu(30000), int(SW * 0.80), Emu(420000), head, 19, INK, bold=True)
        if sub:
            box(s, Emu(1050000), y + Emu(300000), int(SW * 0.80), Emu(520000), sub, 14, GRAY)
            y += Emu(900000)
        else:
            y += Emu(560000)
    if foot:
        box(s, Emu(800000), int(SH * 0.90), int(SW * 0.85), Emu(400000), foot, 12, BRAND)
    return s


# ── 1 封面 ───────────────────────────────────────────────────────────────────
title_slide("AgenticDataHub", "智能实时数据底座 · Intelligent Real-time Data Foundation",
            "以 AI Agent 为操作入口、可直接落地业务的实时数据基础设施 · 本地 Docker 全栈即跑")

# ── 2 产品定位 ───────────────────────────────────────────────────────────────
content_slide(1, "产品定位 · 智能实时数据底座", [
    ("一句话", "把多渠道数据实时归一为 OneID 统一画像，并以 AI Agent 为入口直接圈人、激活、分析。"),
    ("Agentic", "DeepSeek 驱动的自然语言圈人/查询、NL 生成图表与看板、MCP 工具；LLM 绝不直出 SQL（候选 DSL 必过校验层）。"),
    ("DataHub", "小程序/企微/表单/App/批量导入 → 实时 ID-Mapping → OneID → 画像宽表，多租户隔离。"),
], foot="对标 Twilio Segment 的控制台作为应用层：连接 → 用户 → 对象 → 客户 → 触达 → 协议 → 隐私 → 监控 → 知识库 → 应用 → 分析")

# ── 3 架构 ───────────────────────────────────────────────────────────────────
content_slide(2, "整体架构 · 数据链路（接入 → Kafka → Flink → MySQL → Doris）", [
    ("L1 接入层", "多渠道封装为 UserEvent（tenant_id + channel + link_keys + properties）"),
    ("L2 Kafka 事件总线", "大租户独立 Topic / 微租户共享，按 tenant·channel 分区保序"),
    ("L3 Flink 实时计算", "ID-Mapping(OneID 识别/merge) → 画像聚合 → 宽表打宽（dev 用 ID-Mapping 服务 + MySQL 模拟）"),
    ("L3+/L4/L5", "Redis 热层(<5ms) · MySQL 冷层(发号/审计) · Doris OLAP(画像/圈选)"),
], foot="查询层（只读）SQL Engine + 数据 Agent · 调度层 dev 用 Apache Airflow 承载可视化编排")

# ── 4 能力总览 ───────────────────────────────────────────────────────────────
content_slide(3, "核心能力总览", [
    ("实时身份归一", "多渠道 ID-Mapping → OneID，Redis 热层 + MySQL 审计"),
    ("统一画像 + 圈人", "画像宽表；DSL 圈人（多条件/AND·OR/跨对象多跳/边条件）+ 自然语言圈人"),
    ("智能助手 + MCP", "右上角对话助手：DeepSeek + 41 个只读 MCP 工具 + 后台任务"),
    ("连接 / 编排 / 分析", "44 连接器 · Airflow 真实调度 · 知识库 · 应用市场 · NL 分析看板"),
])

# ── 5 ID-Mapping ─────────────────────────────────────────────────────────────
content_slide(4, "实时 ID-Mapping · OneID", [
    ("多渠道识别合并", "wechat openid/unionid、wework extid、phone、email、device 跨渠道归一为 OneID"),
    ("热 + 冷双写", "Redis 写 channel→one_id 热映射（实时点查）；MySQL 写 id_mapping / merge_log 审计"),
    ("可观测", "merge_log 合并审计、画像实时更新；API/Kafka 两种灌入路径"),
], foot="服务 ID-Mapping :8001 · 单文件模拟 Flink ID-Mapping Job")

# ── 6 圈人 ───────────────────────────────────────────────────────────────────
content_slide(5, "统一画像 + 圈人（DSL + 自然语言）", [
    ("统一筛选器", "多条件 + AND/OR；跨对象链式多跳（用户→下单→商品）；关系边条件"),
    ("安全 DSL 层", "候选 DSL 经 validate → echo → compile → estimate；LLM 绝不直出 SQL"),
    ("自然语言圈人", "中文描述 → DeepSeek → 候选 DSL → 实时预估人数 + SQL 预览 → 存为受众"),
], foot="SQL Engine :8002 · 模板 SQL + DSL + NL + 可视化 ETL + 只读 MCP")

# ── 7 智能助手 ───────────────────────────────────────────────────────────────
content_slide(6, "智能助手 · DeepSeek + MCP（右上角常驻）", [
    ("对话即操作", "聊天框里查 schema / 受众 / 对象 / 画像，自然语言驱动"),
    ("桥接 MCP", "stdio 拉起 MCP Server，把 41 个 cdp_* 只读工具桥接成 DeepSeek function-calling"),
    ("发布后台任务", "「发布任务」落到 reverse-ETL 调度模拟后台运行，返回任务状态"),
], foot="独立服务 assistant :8004 · 设置页可查看 MCP 工具清单")

# ── 8 连接器 ─────────────────────────────────────────────────────────────────
content_slide(7, "数据源连接器 · 44 个主流连接器", [
    ("数据库 / NoSQL", "MySQL·PostgreSQL·SQL Server·Oracle·ClickHouse·MariaDB·TiDB·OceanBase·MongoDB·Redis·Cassandra·DynamoDB·Elasticsearch"),
    ("数仓 / 数据湖", "Snowflake·BigQuery·Redshift·Databricks·StarRocks·Doris·Hive · Iceberg·Delta·Hudi·Paimon"),
    ("对象存储 / 流 / 查询", "S3·GCS·ADLS·MinIO·OSS·COS · Kafka·Pulsar·RabbitMQ·Kinesis · Trino·Athena · 国产云数仓 MaxCompute/Hologres/AnalyticDB/DWS"),
], foot="数据源目录页按 8 类平铺；目录级（连接逻辑为占位适配器，可逐个接真实拉数）")

# ── 9 编排 ───────────────────────────────────────────────────────────────────
content_slide(8, "可视化编排 Pipelines · Apache Airflow 真实调度", [
    ("拖拽编排", "左侧连接器节点拖到画布连线 → 保存为管道；管道列表（行表）→ 详情"),
    ("真实触发", "点执行经 Airflow REST API 触发 DAG；参数化通用 DAG 用动态任务映射按节点展开成多任务"),
    ("可观测可控", "详情含拓扑、执行历史、暂停/恢复调度；Airflow 不可达时优雅降级本地模拟"),
], foot="Airflow UI :8088（admin/admin）· DAG: agenticdatahub_pipeline")

# ── 10 知识库 ────────────────────────────────────────────────────────────────
content_slide(9, "知识库 Knowledge Base · 云盘式多模态存储", [
    ("多模态文件", "文档/图片/音视频/压缩包，按类型自动归类 + 图片缩略图；虚拟目录、搜索、类型筛选"),
    ("关联到对象", "文件可关联到对象记录（如 account/A3001），上传时或详情页增删关联"),
    ("真实存储", "文件字节存盘（kb_data 卷），元数据/关联落库；nginx 上传上限放开至 200MB"),
])

# ── 11 应用 + 分析 ───────────────────────────────────────────────────────────
content_slide(10, "应用市场 + 分析看板", [
    ("应用市场", "按类别平铺第三方应用，连接/断开按租户持久化：CRM(Salesforce/HubSpot/销售易)·广告(广点通/巨量/百度)·消息(短信/邮件/企微/钉钉)·分析(神策/GA4)"),
    ("分析 · NL 生成", "一句话描述 → DeepSeek 受限选图 → 直接生成图表/看板；内置用户画像/客户画像/转化率ROI 三看板"),
    ("图表下钻", "点击柱/饼/线的数据点 → 弹出背后明细记录；二次编辑看板（改标题/加减图表）"),
])

# ── 12 技术栈 / 部署 ─────────────────────────────────────────────────────────
content_slide(11, "技术栈 · 部署", [
    ("前端", "React 18 + Vite + TypeScript + TailwindCSS + recharts；对标 Segment 的 IA"),
    ("后端", "FastAPI（SQL Engine / ID-Mapping / 智能助手）+ MySQL + Redis + Kafka + Airflow"),
    ("一键起栈", "docker compose up -d --build（含前端构建）→ http://localhost:8080/console/"),
    ("登录", "强制门禁；演示账号 admin@acme.com / demo123"),
], foot="仓库 github.com/MrRobot2658/agenticdatahub")

# ── 13 结尾 ──────────────────────────────────────────────────────────────────
end = prs.slides.add_slide(BLANK)
bg(end, INK)
rect(end, 0, int(SH * 0.46), SW, Emu(26000), BRAND_LT)
box(end, Emu(700000), int(SH * 0.34), int(SW * 0.85), Emu(900000),
    "AgenticDataHub", 48, WHITE, bold=True)
box(end, Emu(720000), int(SH * 0.52), int(SW * 0.85), Emu(700000),
    "智能实时数据底座 · 让数据被 Agent 直接操作", 22, BRAND_LT)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"已生成 {OUT}（{len(prs.slides.__iter__.__self__._sldIdLst)} 页）")
